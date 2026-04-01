from pptx import Presentation
from pptx.util import Pt


OUT_PATH = r"C:\Users\LENOVO\Desktop\prompt bank\PrePrompt_Full_Project_Presentation.pptx"


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(24)


def main():
    prs = Presentation()

    add_title_slide(
        prs,
        "PrePrompt - Full Project",
        "Unified System: PromptTyper + Prompt Bank\nComplete Final Presentation",
    )

    add_bullets_slide(prs, "Overview + Problem", [
        "PrePrompt = PromptTyper + Prompt Bank as one workflow platform",
        "Prompts scattered in files/notes caused reuse issues",
        "Shortcut workflow lacked centralized prompt source",
        "Manual copy-paste increased errors and wasted time",
        "Need easier UI for non-technical users",
    ])

    add_bullets_slide(prs, "Solution", [
        "Central library with Netflix-style browsing experience",
        "Search, category filters, and quick prompt retrieval",
        "One-click Save to PromptTyper integration",
        "Prompt creation/editing through form, no code changes needed",
    ])

    add_bullets_slide(prs, "Modules + Flow", [
        "PromptTyper (Python): automation + manage prompt shortcuts",
        "Prompt Bank (Web): browse/search/manage prompt library",
        "User browses prompt in Prompt Bank",
        "Clicks Save to PromptTyper",
        "Bridge transfers title/content -> user sets shortcut in PromptTyper",
    ])

    add_bullets_slide(prs, "Architecture Overview", [
        "Frontend Layer: HTML + CSS + JavaScript",
        "Backend/Automation Layer: Python PromptTyper app",
        "Communication Layer: localhost bridge (127.0.0.1:8765)",
        "Storage Layer: local JSON + localStorage metadata",
    ])

    add_bullets_slide(prs, "Database Build Summary", [
        "Prompt records stored in prompt_bank_data.json",
        "Schema includes id, title, category, prompt, images, featured flag",
        "Normalization ensures consistent data format",
        "Pinned/Favorites/Recent saved separately in localStorage",
    ])

    add_bullets_slide(prs, "Bridge API Endpoints", [
        "GET /health -> checks PromptTyper connection status",
        "GET /prompt-bank-data -> loads prompt library",
        "POST /prompt-bank-data -> saves updates from admin module",
        "POST /import-prompt -> sends selected content to PromptTyper",
    ])

    add_bullets_slide(prs, "Core Features", [
        "Netflix-like slider with arrows, dots, and auto slide",
        "Intro loader with controlled sound segment + fade effects",
        "Functional search and quick filters",
        "PromptTyper active/offline status feedback",
        "Add/Edit/Delete prompt admin",
    ])

    add_bullets_slide(prs, "Challenges and Fixes", [
        "No-style page issue fixed with correct asset structure",
        "Slider right-drift bug fixed with stable transform logic",
        "Fetch failures handled with timeout + health checks",
        "Autoplay restrictions handled with tap fallback",
    ])

    add_bullets_slide(prs, "Outcome + Future Scope", [
        "Prompt management became structured and reusable",
        "PromptTyper usage became faster with direct content handoff",
        "System supports both technical and non-technical users",
        "PrePrompt is deployment-ready for local productivity use",
        "Move from JSON to SQLite/PostgreSQL for scale",
        "Cloud sync and user login support",
        "Analytics + team collaboration and role-based access",
    ])

    add_bullets_slide(prs, "Final Conclusion", [
        "PrePrompt = PromptTyper + Prompt Bank as one full product",
        "It solves discovery, management, and execution in one flow",
        "Current version is practical, stable, and extensible",
        "Ready for demo, viva, and team handover",
    ])

    prs.save(OUT_PATH)
    print(f"PPT_CREATED: {OUT_PATH}")


if __name__ == "__main__":
    main()
