from pptx import Presentation
from pptx.util import Pt


OUT_PATH = r"C:\Users\LENOVO\Desktop\prompt bank\PromptBank_Project_Presentation.pptx"


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()

    first = True
    for line in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(24)


def main():
    prs = Presentation()

    add_title_slide(
        prs,
        "Prompt Bank Project",
        "Netflix-style Prompt Library + PromptTyper Integration\nFinal Project Presentation",
    )

    add_bullets_slide(prs, "Problem Statement", [
        "Prompts were scattered in random files/notes",
        "Searching and reusing prompts was slow",
        "Manual code edits needed for each new prompt",
        "PromptTyper integration was not direct",
    ])

    add_bullets_slide(prs, "Our Solution", [
        "Built a centralized Prompt Bank web app",
        "Cinematic browsing with featured slider",
        "Search + category + quick filters for fast discovery",
        "One-click Save to PromptTyper for workflow speed",
    ])

    add_bullets_slide(prs, "Project Objectives", [
        "Make prompt management easy for non-technical users",
        "Allow add/edit/delete without touching code",
        "Connect web prompts to existing Python PromptTyper",
        "Keep system local-first and offline-friendly",
    ])

    add_bullets_slide(prs, "Tech Stack", [
        "Frontend: HTML, CSS, Vanilla JavaScript",
        "Integration: Local HTTP bridge (127.0.0.1:8765)",
        "Storage: prompt_bank_data.json + browser localStorage",
        "Media: Local images + intro audio",
    ])

    add_bullets_slide(prs, "Architecture Overview", [
        "UI Layer -> Home, search, detail, admin panels",
        "Logic Layer -> render, filters, slider, validations",
        "Bridge Layer -> health, load/save, import endpoints",
        "Data Layer -> local JSON records + user meta",
    ])

    add_bullets_slide(prs, "Key Features", [
        "Netflix-style hero carousel (arrows + dots + auto-slide)",
        "Smart search across title/category/tagline/content",
        "Quick filters: All, Pinned, Favorites, Recent",
        "Prompt details: copy + Save to PromptTyper",
    ])

    add_bullets_slide(prs, "Add Prompt / Admin Module", [
        "Mandatory fields: Title + Prompt content",
        "Optional: Image for slider appearance",
        "Edit and Delete from saved prompts list",
        "Poster preview + slider featured toggle",
    ])

    add_bullets_slide(prs, "Database - How We Built It", [
        "Local JSON database model for rapid deployment",
        "Schema: id, title, category, tagline, prompt, images, mainImage, featuredInSlider",
        "Normalization step ensures consistent records",
        "Meta state (pinned/favorite/recent) in localStorage",
    ])

    add_bullets_slide(prs, "API / Bridge Flow", [
        "GET /health -> connection status",
        "GET /prompt-bank-data -> load library",
        "POST /prompt-bank-data -> persist updates",
        "POST /import-prompt -> send selected prompt to PromptTyper",
    ])

    add_bullets_slide(prs, "Intro Experience & Audio", [
        "Intro overlay preloads visuals",
        "Audio plays segment from 5s to 10s",
        "Smooth fade-in and fade-out effect",
        "Fallback: tap to play intro if autoplay blocked",
    ])

    add_bullets_slide(prs, "Challenges We Solved", [
        "No style issue -> fixed HTML/CSS structure",
        "Slider drift bug -> stable transform calculation",
        "Failed fetch errors -> bridge checks + timeout",
        "Layout space waste -> optimized hero alignment",
    ])

    add_bullets_slide(prs, "Testing Checklist", [
        "Slider movement and dot sync validated",
        "Search/filter combinations tested",
        "Add/Edit/Delete prompt CRUD tested",
        "Online/Offline PromptTyper handling tested",
    ])

    add_bullets_slide(prs, "Future Improvements", [
        "Move to SQLite/PostgreSQL for scale",
        "Image storage optimization (avoid heavy base64)",
        "User roles and cloud sync",
        "Analytics dashboard for prompt usage",
    ])

    add_bullets_slide(prs, "Conclusion", [
        "Prompt Bank transforms prompt chaos into structured workflow",
        "Usable for both technical and non-technical users",
        "Direct integration with PromptTyper improves productivity",
        "Project is complete, practical, and ready for extension",
    ])

    prs.save(OUT_PATH)
    print(f"PPT_CREATED: {OUT_PATH}")


if __name__ == "__main__":
    main()
